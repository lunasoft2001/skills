#!/usr/bin/env python3
"""
build_pptx.py — Presentation PPTX Builder
Generates a .pptx file from a storyboard JSON.

Usage:
    python3 build_pptx.py --storyboard /deliverables/my-slug/storyboard.json \
                          --output /deliverables/my-slug/deck.pptx \
                          --theme corporate

Requirements:
    pip install python-pptx
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx is not installed.")
    print("Install it with:  pip install python-pptx")
    sys.exit(1)

# ─── Theme definitions ────────────────────────────────────────────────────────

THEMES = {
    "corporate": {
        "bg":     RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x1B, 0x3A, 0x6B),
        "text":   RGBColor(0x22, 0x22, 0x22),
        "ph_bg":  RGBColor(0xEE, 0xEE, 0xEE),
        "font":   "Calibri",
    },
    "minimal": {
        "bg":     RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x33, 0x33, 0x33),
        "text":   RGBColor(0x33, 0x33, 0x33),
        "ph_bg":  RGBColor(0xF0, 0xF0, 0xF0),
        "font":   "Arial",
    },
    "dark": {
        "bg":     RGBColor(0x1A, 0x1A, 0x2E),
        "accent": RGBColor(0x00, 0xD4, 0xFF),
        "text":   RGBColor(0xF0, 0xF0, 0xF0),
        "ph_bg":  RGBColor(0x2A, 0x2A, 0x3E),
        "font":   "Calibri",
    },
    "vibrant": {
        "bg":     RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x7B, 0x2F, 0xBE),
        "text":   RGBColor(0x1A, 0x1A, 0x1A),
        "ph_bg":  RGBColor(0xF5, 0xF0, 0xFF),
        "font":   "Segoe UI",
    },
}

# ─── Slide dimensions (16:9 widescreen) ──────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── Helpers ─────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height, font_name, font_size, bold=False, color=None, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_placeholder_rect(slide, label, left, top, width, height, theme):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["ph_bg"]
    shape.line.color.rgb = theme["accent"]
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ {label} ]"
    run.font.name = theme["font"]
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = theme["text"]


def add_slide_number(slide, number, theme):
    left  = Inches(12.0)
    top   = Inches(7.1)
    width = Inches(1.0)
    height = Inches(0.3)
    add_text_box(
        slide, str(number), left, top, width, height,
        theme["font"], 11, color=theme["text"], align=PP_ALIGN.RIGHT
    )

# ─── Slide builders ──────────────────────────────────────────────────────────

def build_title_slide(prs, slide_data, theme, project_title):
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    set_bg(slide, theme["bg"])

    # Clear placeholders
    for ph in slide.placeholders:
        ph.element.getparent().remove(ph.element)

    # Accent bar (left side)
    slide.shapes.add_shape(1,
        Inches(0), Inches(0), Inches(0.25), SLIDE_H
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = theme["accent"]
    slide.shapes[-1].line.fill.background()

    # Title
    add_text_box(slide, slide_data.get("title", project_title),
        Inches(0.6), Inches(2.5), Inches(12.0), Inches(1.5),
        theme["font"], 40, bold=True, color=theme["accent"],
        align=PP_ALIGN.LEFT)

    # Subtitle / message
    subtitle = slide_data.get("message", "")
    if subtitle:
        add_text_box(slide, subtitle,
            Inches(0.6), Inches(4.2), Inches(10.0), Inches(1.0),
            theme["font"], 22, color=theme["text"], align=PP_ALIGN.LEFT)

    return slide


def build_content_slide(prs, slide_data, number, theme):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    set_bg(slide, theme["bg"])

    # Accent bar (top)
    bar = slide.shapes.add_shape(1,
        Inches(0), Inches(0), SLIDE_W, Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.fill.background()

    # Title
    add_text_box(slide, slide_data.get("title", ""),
        Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.9),
        theme["font"], 28, bold=True, color=theme["accent"],
        align=PP_ALIGN.LEFT)

    # Message (key takeaway)
    message = slide_data.get("message", "")
    if message:
        add_text_box(slide, message,
            Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.6),
            theme["font"], 16, bold=False, color=theme["text"],
            align=PP_ALIGN.LEFT)

    # Body bullets or visual placeholder
    bullets = slide_data.get("body_bullets", [])
    visual_type = slide_data.get("visual_type", "text")

    needs_placeholder = visual_type not in ("text", "title", "blank", None)

    if bullets and not needs_placeholder:
        body_top = Inches(2.0)
        body_text = "\n".join(f"• {b}" for b in bullets[:5])
        add_text_box(slide, body_text,
            Inches(0.5), body_top, Inches(12.33), Inches(4.5),
            theme["font"], 20, color=theme["text"], align=PP_ALIGN.LEFT)
    elif needs_placeholder:
        label = f"{visual_type.upper()}: {slide_data.get('title', 'Insert visual here')}"
        add_placeholder_rect(slide, label,
            Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5), theme)
        # Show note about placeholder
        add_text_box(slide, "⚠ Replace with actual visual",
            Inches(1.0), Inches(6.6), Inches(5.0), Inches(0.3),
            theme["font"], 10, color=theme["accent"])

    # Slide number
    add_slide_number(slide, number, theme)
    return slide


def build_section_slide(prs, slide_data, number, theme):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, theme["accent"])

    add_text_box(slide, slide_data.get("title", ""),
        Inches(1.0), Inches(2.8), Inches(11.33), Inches(1.5),
        theme["font"], 36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER)

    if slide_data.get("message"):
        add_text_box(slide, slide_data["message"],
            Inches(1.0), Inches(4.5), Inches(11.33), Inches(0.8),
            theme["font"], 18, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER)

    add_slide_number(slide, number, theme)
    return slide

# ─── Main builder ─────────────────────────────────────────────────────────────

def build_presentation(storyboard: dict, output_path: Path, theme_name: str):
    theme = THEMES.get(theme_name, THEMES["corporate"])
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = storyboard.get("slides", [])
    project_title = storyboard.get("title", "Presentation")

    for i, slide_data in enumerate(slides, start=1):
        visual_type = slide_data.get("visual_type", "text")
        act = slide_data.get("act", "")

        if i == 1 or visual_type == "title":
            build_title_slide(prs, slide_data, theme, project_title)
        elif act in ("Section break", "section_break") or visual_type == "section":
            build_section_slide(prs, slide_data, i, theme)
        else:
            build_content_slide(prs, slide_data, i, theme)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"✅ Saved {len(slides)} slides → {output_path}")


def load_example():
    return {
        "title": "Sample Presentation",
        "slug": "sample-deck",
        "slides": [
            {"number": 1, "title": "Sample Presentation", "message": "Built with presentation-pptx-builder", "act": "Title", "visual_type": "title", "body_bullets": []},
            {"number": 2, "title": "The Problem", "message": "Customer churn is costing us $2M per year", "act": "Context", "visual_type": "chart", "body_bullets": []},
            {"number": 3, "title": "Our Solution", "message": "AI-powered retention reduces churn by 40%", "act": "Content", "visual_type": "diagram", "body_bullets": ["Automated risk scoring", "Personalized outreach", "Real-time alerts"]},
            {"number": 4, "title": "Results", "message": "Churn dropped from 8% to 4.5% in Q3", "act": "Content", "visual_type": "chart", "body_bullets": []},
            {"number": 5, "title": "Next Steps", "message": "Three actions to take this quarter", "act": "Action", "visual_type": "text", "body_bullets": ["Approve budget by April 15", "Assign retention team lead", "Set Q3 review checkpoint"]},
        ]
    }


def main():
    parser = argparse.ArgumentParser(description="Build a .pptx from a storyboard JSON")
    parser.add_argument("--storyboard", help="Path to storyboard.json")
    parser.add_argument("--output", default="deck.pptx", help="Output .pptx path")
    parser.add_argument("--theme", default="corporate", choices=list(THEMES.keys()), help="Visual theme")
    parser.add_argument("--example", action="store_true", help="Run with built-in example storyboard")
    args = parser.parse_args()

    if args.example:
        storyboard = load_example()
        output = Path("example_deck.pptx")
    else:
        if not args.storyboard:
            print("ERROR: --storyboard is required (or use --example)")
            sys.exit(1)
        storyboard_path = Path(args.storyboard)
        if not storyboard_path.exists():
            print(f"ERROR: Storyboard not found: {storyboard_path}")
            sys.exit(1)
        with open(storyboard_path, encoding="utf-8") as f:
            storyboard = json.load(f)
        output = Path(args.output)

    build_presentation(storyboard, output, args.theme)


if __name__ == "__main__":
    main()
